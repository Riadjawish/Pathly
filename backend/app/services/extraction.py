"""Text extraction for Pathly's supported study material formats."""

from __future__ import annotations

import asyncio
import re
import zipfile
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from .exceptions import DocumentExtractionError

_MAX_OFFICE_UNCOMPRESSED_BYTES = 250 * 1024 * 1024
_MAX_ZIP_ENTRY_BYTES = 80 * 1024 * 1024


@dataclass(frozen=True, slots=True)
class ExtractedPage:
    number: int
    text: str
    title: str | None = None


@dataclass(frozen=True, slots=True)
class ExtractedDocument:
    text: str
    pages: tuple[ExtractedPage, ...]
    metadata: dict[str, Any] = field(default_factory=dict)


def _normalize_text(value: str) -> str:
    value = value.replace("\x00", " ").replace("\r\n", "\n").replace("\r", "\n")
    value = re.sub(r"[\t\f\v ]+", " ", value)
    value = re.sub(r" *\n *", "\n", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def _join_pages(pages: list[ExtractedPage]) -> str:
    return "\n\n".join(page.text for page in pages if page.text).strip()


def _require_text(document: ExtractedDocument, filename: str) -> ExtractedDocument:
    if not document.text.strip():
        raise DocumentExtractionError(
            f"No selectable text was found in {filename!r}. It may be scanned or image-only."
        )
    return document


def _extract_pdf(path: Path) -> ExtractedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentExtractionError("PDF extraction requires the pypdf package.") from exc

    try:
        reader = PdfReader(str(path), strict=False)
        if reader.is_encrypted:
            try:
                unlocked = reader.decrypt("")
            except Exception as exc:
                raise DocumentExtractionError("Password-protected PDFs are not supported.") from exc
            if not unlocked:
                raise DocumentExtractionError("Password-protected PDFs are not supported.")

        pages: list[ExtractedPage] = []
        for number, page in enumerate(reader.pages, start=1):
            try:
                page_text = _normalize_text(page.extract_text() or "")
            except Exception:
                page_text = ""
            pages.append(ExtractedPage(number=number, text=page_text))

        raw_metadata = reader.metadata or {}
        metadata = {
            str(key).lstrip("/"): str(value)
            for key, value in raw_metadata.items()
            if value is not None
        }
        metadata["page_count"] = len(pages)
        return _require_text(
            ExtractedDocument(text=_join_pages(pages), pages=tuple(pages), metadata=metadata),
            path.name,
        )
    except DocumentExtractionError:
        raise
    except Exception as exc:
        raise DocumentExtractionError(f"Could not read PDF {path.name!r}.") from exc


def _validate_office_archive(path: Path, expected_directory: str) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if "[Content_Types].xml" not in names or not any(
                name.startswith(f"{expected_directory}/") for name in names
            ):
                raise DocumentExtractionError(
                    "The Office document has an invalid internal structure."
                )

            total = 0
            for info in archive.infolist():
                if info.file_size > _MAX_ZIP_ENTRY_BYTES:
                    raise DocumentExtractionError(
                        "The Office document contains an oversized entry."
                    )
                total += info.file_size
                if total > _MAX_OFFICE_UNCOMPRESSED_BYTES:
                    raise DocumentExtractionError(
                        "The Office document expands beyond the safe size limit."
                    )
    except zipfile.BadZipFile as exc:
        raise DocumentExtractionError("The Office document is not a valid ZIP archive.") from exc


def _extract_docx(path: Path) -> ExtractedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentExtractionError("Word extraction requires the python-docx package.") from exc

    _validate_office_archive(path, "word")
    try:
        document = Document(str(path))
        blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                values = [_normalize_text(cell.text) for cell in row.cells]
                if any(values):
                    blocks.append(" | ".join(values))
        text = _normalize_text("\n\n".join(blocks))
        page = ExtractedPage(number=1, text=text)
        metadata = {
            "paragraph_count": len(document.paragraphs),
            "table_count": len(document.tables),
        }
        return _require_text(
            ExtractedDocument(text=text, pages=(page,), metadata=metadata), path.name
        )
    except DocumentExtractionError:
        raise
    except Exception as exc:
        raise DocumentExtractionError(f"Could not read Word document {path.name!r}.") from exc


def _extract_pptx(path: Path) -> ExtractedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentExtractionError(
            "PowerPoint extraction requires the python-pptx package."
        ) from exc

    _validate_office_archive(path, "ppt")
    try:
        presentation = Presentation(str(path))
        pages: list[ExtractedPage] = []
        for number, slide in enumerate(presentation.slides, start=1):
            blocks: list[str] = []
            title: str | None = None
            if slide.shapes.title is not None:
                title = _normalize_text(slide.shapes.title.text or "") or None

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = _normalize_text(getattr(shape, "text", "") or "")
                    if text:
                        blocks.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [_normalize_text(cell.text) for cell in row.cells]
                        if any(values):
                            blocks.append(" | ".join(values))

            try:
                notes = slide.notes_slide.notes_text_frame.text
            except Exception:
                notes = ""
            if notes and _normalize_text(notes) not in blocks:
                blocks.append(f"Speaker notes:\n{_normalize_text(notes)}")

            pages.append(
                ExtractedPage(number=number, text=_normalize_text("\n\n".join(blocks)), title=title)
            )

        return _require_text(
            ExtractedDocument(
                text=_join_pages(pages),
                pages=tuple(pages),
                metadata={"slide_count": len(pages)},
            ),
            path.name,
        )
    except DocumentExtractionError:
        raise
    except Exception as exc:
        raise DocumentExtractionError(f"Could not read presentation {path.name!r}.") from exc


def _extract_text(path: Path) -> ExtractedDocument:
    try:
        text = _normalize_text(path.read_text(encoding="utf-8-sig"))
    except UnicodeDecodeError as exc:
        raise DocumentExtractionError("Text documents must use UTF-8 encoding.") from exc
    except OSError as exc:
        raise DocumentExtractionError(f"Could not read text document {path.name!r}.") from exc
    page = ExtractedPage(number=1, text=text)
    return _require_text(
        ExtractedDocument(text=text, pages=(page,), metadata={"character_count": len(text)}),
        path.name,
    )


_EXTRACTORS: dict[str, Callable[[Path], ExtractedDocument]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".txt": _extract_text,
    ".md": _extract_text,
}


async def extract_document(path: str | Path) -> ExtractedDocument:
    """Extract normalized text without blocking the API event loop."""

    local_path = Path(path)
    extractor = _EXTRACTORS.get(local_path.suffix.lower())
    if extractor is None:
        supported = ", ".join(sorted(_EXTRACTORS))
        raise DocumentExtractionError(f"Unsupported document type. Supported types: {supported}.")
    if not await asyncio.to_thread(local_path.is_file):
        raise DocumentExtractionError("The document file does not exist.")
    return await asyncio.to_thread(extractor, local_path)
